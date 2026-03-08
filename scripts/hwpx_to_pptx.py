"""
HWPX 제안서 → PPTX 발표자료 변환 스크립트 v2
2단계 Claude 파이프라인 + 10종 레이아웃 + 스토리보드 최적화

사용법: uv run python scripts/hwpx_to_pptx.py [hwpx_file]
"""
import zipfile
import json
import sys
import os
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import anthropic
from dotenv import load_dotenv

load_dotenv()
sys.stdout.reconfigure(encoding="utf-8")

# ── 색상 팔레트 ──────────────────────────────────────────
PRIMARY    = RGBColor(0x1A, 0x3A, 0x6B)   # 진청 (헤더, 표지)
SECONDARY  = RGBColor(0x00, 0x8C, 0xCB)   # 하늘청 (포인트)
ACCENT     = RGBColor(0xFF, 0x8C, 0x00)   # 주황 (강조)
ACCENT2    = RGBColor(0x00, 0xB0, 0x5A)   # 초록 (긍정 지표)
BG_LIGHT   = RGBColor(0xF4, 0xF7, 0xFB)   # 연한 배경
BG_CARD    = RGBColor(0xFF, 0xFF, 0xFF)   # 카드 배경
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_TEXT  = RGBColor(0x66, 0x66, 0x77)
LIGHT_LINE = RGBColor(0xDD, 0xE5, 0xF0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

CARD_COLORS = [PRIMARY, SECONDARY, ACCENT, ACCENT2]


# ════════════════════════════════════════════════════════
# 1. HWPX 텍스트 추출
# ════════════════════════════════════════════════════════
def extract_hwpx_text(hwpx_path: str) -> str:
    HP = "http://www.hancom.co.kr/hwpml/2011/paragraph"
    with zipfile.ZipFile(hwpx_path, "r") as z:
        xml_data = z.read("Contents/section0.xml")
    root = etree.fromstring(xml_data)
    texts = [t.text for t in root.iter(f"{{{HP}}}t") if t.text]
    return "".join(texts)


# ════════════════════════════════════════════════════════
# 2. Step 1 — 문서 핵심 정보 구조화 추출
# ════════════════════════════════════════════════════════
STEP1_PROMPT = """당신은 공공기관 용역 입찰 제안 PT 전문가입니다.
아래 제안서를 분석하여 평가위원을 설득하는 발표용 PPT 스토리보드를 설계하세요.

## 배경 지식 (반드시 반영)

### 평가 구조
- 협상에 의한 계약: 기술능력 70점 + 가격 30점
- 승패는 PT Q&A에서 결정됨 → 각 슬라이드에 예상 질문 포함 필수

### 이기는 PT의 스토리 흐름 (반드시 이 순서를 따름)
```
1섹션 문제 공감  → 발주처 상황을 우리가 가장 잘 이해함을 증명
2섹션 솔루션     → 일반 접근법 대비 우리 방법의 우월성 (two_column 활용)
3섹션 실행 역량  → 수치+레퍼런스로 실행 가능성 검증
4섹션 성공 확신  → 기대효과, 추진일정, 마무리
```

### 슬라이드 설계 원칙 (One Slide, One Message)
- 슬라이드당 핵심 메시지 1개만
- 6×6 Rule: bullet 최대 6개, 각 bullet 최대 6단어
- F-Pattern: 숫자·다이어그램 → 제목 → 본문 순으로 시선 유도
- 숫자/통계는 반드시 numbers_callout 레이아웃으로 분리

### two_column 필수 활용 패턴
- left_title: "일반적 접근법" / right_title: "우리의 방법론"
- 경쟁사와 직접 비교 금지, 대신 "일반 방법 vs 우리 방법"으로 표현

---

제안서 텍스트:
{text}

---

아래 JSON 형식으로만 출력하세요:
{{
  "title": "제안서 제목",
  "subtitle": "한 줄 부제목",
  "company": "제안사명",
  "date": "제출일",
  "client": "발주처명",
  "project_summary": "사업 목적 2~3문장 요약",
  "sections": [
    {{
      "section_num": "01",
      "section_title": "섹션명",
      "slides": [
        {{
          "purpose": "이 슬라이드의 목적 한 줄",
          "layout": "layout_type",
          "title": "슬라이드 제목",
          "key_message": "평가위원이 기억해야 할 핵심 결론 (15자 이내)",
          "content": {{}},
          "speaker_notes": "발표자 가이드: 강조 포인트 + 예상 Q&A 1~2개"
        }}
      ]
    }}
  ],
  "closing_message": "마무리 메시지"
}}

layout_type 선택 기준:
- "agenda"          : 목차 (표지 직후 1개)
- "numbers_callout" : 숫자/통계 3~4개 강조 — 임팩트 데이터 있을 때 필수
- "process_flow"    : 3~5단계 순서 — 프로세스/방법론 설명
- "two_column"      : 비교/대조 — 반드시 "일반 vs 우리" 패턴으로 사용
- "timeline"        : 추진일정 — 월별/단계별
- "quote_highlight" : 비전/미션/슬로건 강조
- "content"         : 위 패턴이 맞지 않는 일반 내용
- "closing"         : 마무리

content 필드 형식 (layout_type별):
- agenda:           {{"sections": [{{"num":"01","title":"섹션명"}}]}}
- numbers_callout:  {{"numbers": [{{"value":"89%","label":"6단어 이내 설명"}}]}}
- process_flow:     {{"steps": [{{"num":"1","title":"동사형 단계명","desc":"20자 이내"}}]}}
- two_column:       {{"left_title":"일반적 접근법","left_bullets":["6단어 이내"],"right_title":"우리의 방법론","right_bullets":["6단어 이내"]}}
- timeline:         {{"items": [{{"period":"1~2월","title":"단계명","tasks":["세부 1개"]}}]}}
- quote_highlight:  {{"quote":"핵심 문장","attribution":"근거/출처"}}
- content:          {{"bullets": ["6단어 이내 항목"]}}
- closing:          {{"message":"감사 메시지","cta":"Q & A"}}

제약 조건:
- sections 최대 4개, 각 section 내 slides 2~4장
- 총 콘텐츠 슬라이드 12~16장 (section_header 미포함)
- 같은 layout 3장 연속 금지
- numbers_callout 최소 1개, process_flow 최소 1개, two_column 최소 1개
- 모든 bullet/desc/label/tasks 항목: 6단어(30자) 이내
- key_message: 15자 이내 결론형 문장
- speaker_notes: 발표 포인트 + "예상Q: ~" 형식으로 1~2개 포함"""


# ════════════════════════════════════════════════════════
# 3. Step 2 — 슬라이드 콘텐츠 정제 (선택적)
# ════════════════════════════════════════════════════════
STEP2_PROMPT = """당신은 공공 입찰 PT 전문가입니다. 다음 슬라이드 구조를 검토하고 수주 확률을 높이도록 개선하세요.

현재 구조:
{structure}

## 체크리스트 (모두 통과해야 함)

### 스토리 구조
- [ ] 섹션 순서가 "문제 공감 → 솔루션 → 실행 역량 → 성공 확신" 흐름인가?
- [ ] 첫 섹션 첫 슬라이드가 발주처 문제를 구체적 수치로 제시하는가?
- [ ] two_column이 "일반적 접근법 vs 우리의 방법론" 패턴으로 사용되는가?

### 슬라이드 품질 (One Slide, One Message)
- [ ] key_message가 15자 이내 결론형 문장인가? (예: "20년 노하우로 확신")
- [ ] bullet이 6단어 이내인가? (예: "전문가 DB 6만명 보유" ✅ / "저희 회사는 20년간 다양한 분야에서 전문가 네트워크를 구축해왔습니다" ❌)
- [ ] numbers_callout의 value가 임팩트 있는 단위인가? ("89%" "3배↑" "6만명" "20년")

### Q&A 대비 (speaker_notes)
- [ ] 각 슬라이드에 "예상Q:" 형식 질문이 1~2개 포함되어 있는가?
- [ ] 취약점 슬라이드(일정, 예산, 인력)에 방어 포인트가 명시되어 있는가?

### 시각 계층 (F-Pattern)
- [ ] process_flow의 step title이 동사형인가? ("테마 발굴" ✅ / "테마" ❌)
- [ ] numbers_callout의 label이 6단어 이내인가?

위 체크리스트를 기준으로 미통과 항목을 모두 수정하세요.
동일한 JSON 구조로 개선된 내용만 출력하세요."""


def call_claude(prompt: str, max_tokens: int = 8000) -> str:
    client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    msg = client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = msg.content[0].text.strip()
    if "```json" in raw:
        raw = raw.split("```json")[1].split("```")[0].strip()
    elif "```" in raw:
        raw = raw.split("```")[1].split("```")[0].strip()
    return raw


def generate_storyboard(text: str) -> dict:
    print("  [Step 1] 문서 구조 분석 및 스토리보드 설계...")
    raw1 = call_claude(STEP1_PROMPT.format(text=text[:28000]))
    doc = json.loads(raw1)

    print("  [Step 2] 슬라이드 콘텐츠 정제...")
    raw2 = call_claude(STEP2_PROMPT.format(structure=json.dumps(doc, ensure_ascii=False, indent=2)))
    doc = json.loads(raw2)

    return doc


# ════════════════════════════════════════════════════════
# 4. PPTX 공통 유틸
# ════════════════════════════════════════════════════════
def set_speaker_notes(slide, notes: str):
    """발표자 노트(Q&A 예상 질문 포함)를 슬라이드에 저장"""
    if not notes:
        return
    from pptx.util import Pt
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes


def add_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color: RGBColor, line_color=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    return shp


def txb(slide, text, l, t, w, h,
        size=16, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def header_bar(slide, title: str, slide_num: int = 0):
    """공통 헤더 바 + 제목 + 슬라이드 번호"""
    rect(slide, 0, 0, SLIDE_W, Inches(1.1), PRIMARY)
    rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.06), SECONDARY)
    txb(slide, title,
        Inches(0.55), Inches(0.12), Inches(12.0), Inches(0.9),
        size=26, bold=True, color=WHITE)
    if slide_num:
        txb(slide, str(slide_num),
            Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3),
            size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


def key_msg_bar(slide, msg: str):
    """핵심 메시지 강조 바"""
    rect(slide, Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.65),
         RGBColor(0xE6, 0xF2, 0xFB))
    rect(slide, Inches(0.5), Inches(1.3), Inches(0.06), Inches(0.65), SECONDARY)
    txb(slide, msg,
        Inches(0.68), Inches(1.33), Inches(12.0), Inches(0.58),
        size=15, bold=True, color=PRIMARY)


def bullet_list(slide, bullets: list, start_y, left=Inches(0.55), width=Inches(12.3)):
    y = start_y
    for pt in bullets[:5]:
        rect(slide, left, y + Inches(0.13), Inches(0.1), Inches(0.3), SECONDARY)
        txb(slide, pt, left + Inches(0.2), y, width - Inches(0.2), Inches(0.55),
            size=16, color=DARK_TEXT)
        y += Inches(0.65)
    return y


# ════════════════════════════════════════════════════════
# 5. 슬라이드 빌더 — 레이아웃별
# ════════════════════════════════════════════════════════
def build_cover(prs, doc: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    # 배경 장식
    rect(slide, 0, 0, SLIDE_W, Inches(0.1), SECONDARY)
    rect(slide, 0, Inches(7.4), SLIDE_W, Inches(0.1), ACCENT)
    rect(slide, Inches(11.3), Inches(0.1), Inches(2.03), Inches(7.3),
         RGBColor(0x14, 0x2F, 0x5E))
    rect(slide, Inches(11.3), Inches(0.1), Inches(0.06), Inches(7.3), SECONDARY)

    # 클라이언트 라벨
    client = doc.get("client", "")
    if client:
        txb(slide, f"제출처: {client}",
            Inches(0.7), Inches(1.5), Inches(10.0), Inches(0.5),
            size=14, color=RGBColor(0x88, 0xAA, 0xCC))

    # 메인 제목
    txb(slide, doc.get("title", ""),
        Inches(0.7), Inches(2.1), Inches(10.2), Inches(2.0),
        size=38, bold=True, color=WHITE)

    # 부제목
    subtitle = doc.get("subtitle", "")
    if subtitle:
        txb(slide, subtitle,
            Inches(0.7), Inches(4.0), Inches(10.2), Inches(0.8),
            size=20, color=SECONDARY)

    # 구분선
    rect(slide, Inches(0.7), Inches(4.85), Inches(1.5), Inches(0.04), ACCENT)

    # 날짜 & 회사
    footer = "   |   ".join(filter(None, [doc.get("date", ""), doc.get("company", "")]))
    txb(slide, footer,
        Inches(0.7), Inches(5.1), Inches(10.2), Inches(0.5),
        size=15, color=RGBColor(0xBB, 0xCC, 0xDD))


def build_section_header(prs, section_num: str, section_title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x0D, 0x27, 0x55))
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), ACCENT)
    rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), SECONDARY)

    # 배경 장식 원형 효과 (큰 사각형으로 근사)
    rect(slide, Inches(9.5), Inches(1.0), Inches(5.0), Inches(5.5),
         RGBColor(0x12, 0x35, 0x65))

    txb(slide, section_num,
        Inches(1.2), Inches(2.2), Inches(3.5), Inches(2.0),
        size=80, bold=True, color=RGBColor(0x1A, 0x5A, 0xAA), align=PP_ALIGN.LEFT)
    rect(slide, Inches(1.2), Inches(4.4), Inches(0.07), Inches(1.1), ACCENT)
    txb(slide, section_title,
        Inches(1.5), Inches(4.3), Inches(10.0), Inches(1.3),
        size=34, bold=True, color=WHITE)


def build_agenda(prs, sections: list, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_LIGHT)
    header_bar(slide, "목  차", slide_num)
    rect(slide, 0, Inches(1.16), SLIDE_W, Inches(0.06), SECONDARY)

    cols = 2
    items_per_col = (len(sections) + 1) // 2
    col_w = Inches(5.8)
    for i, sec in enumerate(sections):
        col = i // items_per_col
        row = i % items_per_col
        lx = Inches(0.8) + col * Inches(6.7)
        ty = Inches(1.7) + row * Inches(0.88)

        # 번호 배지
        rect(slide, lx, ty, Inches(0.55), Inches(0.55),
             CARD_COLORS[i % len(CARD_COLORS)])
        txb(slide, sec.get("num", str(i + 1)),
            lx, ty, Inches(0.55), Inches(0.55),
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # 제목
        txb(slide, sec.get("title", ""),
            lx + Inches(0.65), ty + Inches(0.05), col_w - Inches(0.8), Inches(0.5),
            size=18, bold=True, color=DARK_TEXT)

        # 구분선
        rect(slide, lx + Inches(0.65), ty + Inches(0.58), col_w - Inches(1.0),
             Inches(0.02), LIGHT_LINE)


def build_content(prs, data: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_LIGHT)
    header_bar(slide, data.get("title", ""), slide_num)

    key_msg = data.get("key_message", "")
    if key_msg:
        key_msg_bar(slide, key_msg)
        start_y = Inches(2.1)
    else:
        start_y = Inches(1.35)

    bullets = data.get("content", {}).get("bullets", [])
    bullet_list(slide, bullets, start_y)


def build_two_column(prs, data: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_LIGHT)
    header_bar(slide, data.get("title", ""), slide_num)

    key_msg = data.get("key_message", "")
    if key_msg:
        key_msg_bar(slide, key_msg)
        col_top = Inches(2.1)
    else:
        col_top = Inches(1.35)

    c = data.get("content", {})
    col_h = Inches(7.3) - col_top - Inches(0.3)

    for i, (col_color, title_key, bullets_key, lx) in enumerate([
        (PRIMARY,   "left_title",  "left_bullets",  Inches(0.5)),
        (SECONDARY, "right_title", "right_bullets", Inches(6.95)),
    ]):
        # 컬럼 배경 카드
        rect(slide, lx, col_top, Inches(6.0), col_h,
             BG_CARD, LIGHT_LINE)
        # 헤더
        rect(slide, lx, col_top, Inches(6.0), Inches(0.55), col_color)
        txb(slide, c.get(title_key, ""),
            lx + Inches(0.15), col_top + Inches(0.06),
            Inches(5.7), Inches(0.45),
            size=15, bold=True, color=WHITE)
        # 불릿
        bullet_list(slide, c.get(bullets_key, []),
                    col_top + Inches(0.65), lx + Inches(0.2), Inches(5.6))


def build_numbers_callout(prs, data: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_LIGHT)
    header_bar(slide, data.get("title", ""), slide_num)

    key_msg = data.get("key_message", "")
    if key_msg:
        key_msg_bar(slide, key_msg)
        card_top = Inches(2.15)
    else:
        card_top = Inches(1.4)

    numbers = data.get("content", {}).get("numbers", [])[:4]
    n = len(numbers)
    card_w = Inches(12.3 / n) - Inches(0.2)
    card_h = SLIDE_H - card_top - Inches(0.5)
    gap = Inches(0.25)

    for i, num in enumerate(numbers):
        lx = Inches(0.5) + i * (card_w + gap)
        color = CARD_COLORS[i % len(CARD_COLORS)]

        # 카드 배경
        rect(slide, lx, card_top, card_w, card_h, BG_CARD, LIGHT_LINE)
        # 상단 컬러 바
        rect(slide, lx, card_top, card_w, Inches(0.12), color)

        # 큰 숫자
        txb(slide, num.get("value", ""),
            lx + Inches(0.1), card_top + Inches(0.3),
            card_w - Inches(0.2), Inches(1.5),
            size=48, bold=True, color=color, align=PP_ALIGN.CENTER)

        # 라벨
        txb(slide, num.get("label", ""),
            lx + Inches(0.1), card_top + Inches(1.9),
            card_w - Inches(0.2), Inches(1.5),
            size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


def build_process_flow(prs, data: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_LIGHT)
    header_bar(slide, data.get("title", ""), slide_num)

    key_msg = data.get("key_message", "")
    if key_msg:
        key_msg_bar(slide, key_msg)
        proc_top = Inches(2.1)
    else:
        proc_top = Inches(1.4)

    steps = data.get("content", {}).get("steps", [])[:5]
    n = len(steps)
    if n == 0:
        return

    total_w = Inches(12.3)
    step_w = total_w / n - Inches(0.1)
    step_h = Inches(3.5)
    arrow_w = Inches(0.35)

    for i, step in enumerate(steps):
        lx = Inches(0.5) + i * (step_w + arrow_w)
        color = CARD_COLORS[i % len(CARD_COLORS)]

        # 스텝 박스
        rect(slide, lx, proc_top, step_w, step_h, BG_CARD, LIGHT_LINE)
        rect(slide, lx, proc_top, step_w, Inches(0.55), color)

        # 번호
        txb(slide, step.get("num", str(i + 1)),
            lx + Inches(0.1), proc_top + Inches(0.05),
            Inches(0.45), Inches(0.45),
            size=22, bold=True, color=WHITE)

        # 단계 제목
        txb(slide, step.get("title", ""),
            lx + Inches(0.55), proc_top + Inches(0.06),
            step_w - Inches(0.65), Inches(0.45),
            size=14, bold=True, color=WHITE)

        # 설명
        txb(slide, step.get("desc", ""),
            lx + Inches(0.12), proc_top + Inches(0.68),
            step_w - Inches(0.24), step_h - Inches(0.8),
            size=13, color=DARK_TEXT, wrap=True)

        # 화살표 (마지막 제외)
        if i < n - 1:
            ax = lx + step_w + Inches(0.05)
            ay = proc_top + step_h / 2 - Inches(0.15)
            txb(slide, "▶",
                ax, ay, arrow_w, Inches(0.3),
                size=18, color=SECONDARY, align=PP_ALIGN.CENTER)


def build_timeline(prs, data: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_LIGHT)
    header_bar(slide, data.get("title", ""), slide_num)

    key_msg = data.get("key_message", "")
    if key_msg:
        key_msg_bar(slide, key_msg)
        tl_top = Inches(2.15)
    else:
        tl_top = Inches(1.4)

    items = data.get("content", {}).get("items", [])[:6]
    n = len(items)
    if n == 0:
        return

    col_w = Inches(12.3) / n
    line_y = tl_top + Inches(0.85)

    # 중앙 타임라인 선
    rect(slide, Inches(0.5), line_y, Inches(12.3), Inches(0.06), SECONDARY)

    for i, item in enumerate(items):
        lx = Inches(0.5) + i * col_w
        color = CARD_COLORS[i % len(CARD_COLORS)]

        # 타임라인 노드
        node_x = lx + col_w / 2 - Inches(0.2)
        rect(slide, node_x, line_y - Inches(0.2), Inches(0.4), Inches(0.4), color)

        # 기간 라벨 (노드 위)
        txb(slide, item.get("period", ""),
            lx, tl_top, col_w, Inches(0.7),
            size=13, bold=True, color=color, align=PP_ALIGN.CENTER)

        # 제목 (노드 아래)
        txb(slide, item.get("title", ""),
            lx + Inches(0.05), line_y + Inches(0.35),
            col_w - Inches(0.1), Inches(0.6),
            size=14, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

        # 세부 태스크
        tasks = item.get("tasks", [])
        ty = line_y + Inches(1.05)
        for task in tasks[:3]:
            txb(slide, f"• {task}",
                lx + Inches(0.08), ty,
                col_w - Inches(0.16), Inches(0.5),
                size=12, color=GRAY_TEXT, align=PP_ALIGN.LEFT)
            ty += Inches(0.48)


def build_quote_highlight(prs, data: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), ACCENT)
    rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), SECONDARY)

    title = data.get("title", "")
    if title:
        txb(slide, title,
            Inches(1.0), Inches(0.3), Inches(11.0), Inches(0.7),
            size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

    c = data.get("content", {})
    quote = c.get("quote", "")
    attribution = c.get("attribution", "")

    # 큰따옴표 장식
    txb(slide, "\u201c",
        Inches(0.8), Inches(1.5), Inches(1.5), Inches(2.0),
        size=120, color=SECONDARY, align=PP_ALIGN.LEFT)

    txb(slide, quote,
        Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5),
        size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    if attribution:
        rect(slide, Inches(5.5), Inches(5.0), Inches(2.33), Inches(0.04), ACCENT)
        txb(slide, attribution,
            Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6),
            size=16, color=SECONDARY, align=PP_ALIGN.CENTER)

    if slide_num:
        txb(slide, str(slide_num),
            Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3),
            size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


def build_closing(prs, data: dict, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    rect(slide, 0, 0, SLIDE_W, Inches(0.1), ACCENT)
    rect(slide, 0, Inches(7.4), SLIDE_W, Inches(0.1), SECONDARY)

    # 배경 장식
    rect(slide, Inches(9.0), Inches(1.0), Inches(5.5), Inches(5.5),
         RGBColor(0x14, 0x2F, 0x5E))

    c = data.get("content", {})
    msg = c.get("message", "감사합니다")
    cta = c.get("cta", "Q & A")

    txb(slide, msg,
        Inches(1.0), Inches(2.2), Inches(9.0), Inches(1.8),
        size=42, bold=True, color=WHITE)
    rect(slide, Inches(1.0), Inches(4.1), Inches(2.5), Inches(0.06), ACCENT)
    txb(slide, cta,
        Inches(1.0), Inches(4.3), Inches(9.0), Inches(0.8),
        size=24, color=SECONDARY)

    # 회사 정보
    company = data.get("company", "")
    if company:
        txb(slide, company,
            Inches(1.0), Inches(6.5), Inches(9.0), Inches(0.5),
            size=14, color=RGBColor(0x88, 0xAA, 0xCC))


# ════════════════════════════════════════════════════════
# 6. 슬라이드 디스패처
# ════════════════════════════════════════════════════════
LAYOUT_BUILDERS = {
    "content":          build_content,
    "two_column":       build_two_column,
    "numbers_callout":  build_numbers_callout,
    "process_flow":     build_process_flow,
    "timeline":         build_timeline,
    "quote_highlight":  build_quote_highlight,
}


def build_pptx(doc: dict, output_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_num = 0

    # 1. 표지
    build_cover(prs, doc)
    slide_num += 1

    # 2. 목차 (전체 섹션 수집)
    sections_meta = [
        {"num": s.get("section_num", str(i+1)), "title": s.get("section_title", "")}
        for i, s in enumerate(doc.get("sections", []))
    ]
    slide_num += 1
    build_agenda(prs, sections_meta, slide_num)

    # 3. 각 섹션
    for section in doc.get("sections", []):
        # 섹션 헤더
        build_section_header(prs,
                             section.get("section_num", ""),
                             section.get("section_title", ""))

        # 섹션 내 슬라이드
        for slide_data in section.get("slides", []):
            layout = slide_data.get("layout", "content")
            if layout == "agenda":
                continue  # 이미 앞에서 처리

            slide_num += 1
            builder = LAYOUT_BUILDERS.get(layout, build_content)
            # closing은 마지막에 따로 처리
            if layout == "closing":
                slide_data["company"] = doc.get("company", "")
                build_closing(prs, slide_data, slide_num)
            else:
                builder(prs, slide_data, slide_num)

            # 발표자 노트 저장 (Q&A 예상 질문 포함)
            notes = slide_data.get("speaker_notes", "")
            if notes:
                set_speaker_notes(prs.slides[-1], notes)

    # 4. 마지막 closing이 없으면 추가
    last_layouts = [
        s.get("layout") for sec in doc.get("sections", [])
        for s in sec.get("slides", [])
    ]
    if "closing" not in last_layouts:
        slide_num += 1
        build_closing(prs, {
            "content": {
                "message": "감사합니다",
                "cta": "Q & A"
            },
            "company": doc.get("company", "")
        }, slide_num)

    prs.save(output_path)
    print(f"saved: {output_path}  ({len(prs.slides)} slides)")


# ════════════════════════════════════════════════════════
# 7. 메인
# ════════════════════════════════════════════════════════
def main():
    hwpx_path = sys.argv[1] if len(sys.argv) > 1 else \
        "output/(제안서 샘플) O-Prize 사업기획 및 운영 용역_사업제안서.hwpx"

    print(f"[1/3] HWPX parsing: {hwpx_path}")
    text = extract_hwpx_text(hwpx_path)
    print(f"      {len(text):,} chars extracted")

    print("[2/3] Storyboard generation (2-step Claude pipeline)...")
    doc = generate_storyboard(text)

    sections = doc.get("sections", [])
    total_slides = sum(len(s.get("slides", [])) for s in sections)
    print(f"      {len(sections)} sections, {total_slides} content slides")

    stem = Path(hwpx_path).stem
    output_path = f"output/{stem}_발표자료_v2.pptx"
    print(f"[3/3] Building PPTX...")
    build_pptx(doc, output_path)


if __name__ == "__main__":
    main()
