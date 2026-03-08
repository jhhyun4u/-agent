"""발표 자료 전용 PPTX 빌더 — 컨설팅급 시각 설계 원칙 반영

설계 원칙:
- F-Pattern 시선 흐름: 큰 숫자/차트 → 제목 → 본문
- 하이라이트 1개 원칙: 슬라이드당 bold+색상 강조 요소 1개
- 차트 제목 = 결론 문장
- 레이아웃 14종: cover/agenda/key_message/eval_section/comparison/timeline/team/closing/
                  numbers_callout/process_flow/problem_sync/quote_highlight/
                  split_panel/numbered_strategy/case_study
"""

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# ── 색상 상수 ────────────────────────────────────────────────────────────────
COLOR_DARK_BLUE  = RGBColor(0x1F, 0x49, 0x7D)   # 제목, eval_badge, 강조
COLOR_ACCENT     = RGBColor(0x2E, 0x75, 0xB6)   # 소제목, 구분선, highlight
COLOR_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)   # 표 짝수행 배경
COLOR_DARK_TEXT  = RGBColor(0x26, 0x26, 0x26)   # 본문
COLOR_SUB_TEXT   = RGBColor(0x70, 0x70, 0x70)   # sub-bullet, 보조 설명
COLOR_OUR_BG     = RGBColor(0xD6, 0xE4, 0xF0)   # comparison "우리" 열 배경
COLOR_ACCENT_BAR = RGBColor(0x1F, 0x49, 0x7D)   # 제목 좌측 accent bar
COLOR_TITLE_BG   = RGBColor(0xEA, 0xF1, 0xFB)   # 제목 배경 밴드
COLOR_YELLOW     = RGBColor(0xFF, 0xC0, 0x00)   # problem_sync 수치 강조
COLOR_CASE_BG    = RGBColor(0xF7, 0xFA, 0xFD)   # case_study 카드 배경

# ── 폰트 ──────────────────────────────────────────────────────────────────────
FONT_NAME = "맑은 고딕"


# ── 유틸 헬퍼 ─────────────────────────────────────────────────────────────────

def _add_textbox(slide, left, top, width, height, text="", font_size=18,
                 bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    """텍스트박스 추가 헬퍼"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def _add_eval_badge(slide, badge_text: str):
    """슬라이드 우측 상단에 [평가항목명 | XX점] 배지 표시"""
    if not badge_text:
        return
    box = slide.shapes.add_shape(
        1,
        Inches(9.3), Inches(0.1), Inches(3.8), Inches(0.42),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEA, 0xF1, 0xFB)
    box.line.color.rgb = COLOR_ACCENT
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = badge_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = FONT_NAME
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = COLOR_DARK_BLUE


def _add_speaker_notes(slide, notes_text: str):
    """발표자 노트 추가"""
    if not notes_text:
        return
    slide.notes_slide.notes_text_frame.text = notes_text


def _add_slide_number(slide, num: int):
    """슬라이드 우측 하단에 페이지 번호 표시 (표지 제외)"""
    _add_textbox(
        slide,
        Inches(12.0), Inches(6.9), Inches(1.1), Inches(0.35),
        text=str(num),
        font_size=12,
        color=RGBColor(0x80, 0x80, 0x80),
        align=PP_ALIGN.RIGHT,
    )


def _add_slide_title(slide, title: str, top=Inches(0.15)):
    """슬라이드 제목 — 좌측 accent bar + 연한 배경 밴드 + 28pt bold

    기존 얇은 구분선 대신 컨설팅급 제목 영역 처리:
    - 슬라이드 전체 폭 연한 배경 밴드
    - 좌측 4pt 굵기 강조 bar
    - 28pt bold 제목
    """
    # 배경 밴드 (슬라이드 전체 폭)
    band = slide.shapes.add_shape(
        1,
        Inches(0), top, Inches(13.333), Inches(0.72),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_TITLE_BG
    band.line.fill.background()

    # 좌측 accent bar
    bar = slide.shapes.add_shape(
        1,
        Inches(0), top, Inches(0.08), Inches(0.72),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT_BAR
    bar.line.fill.background()

    # 제목 텍스트
    _add_textbox(
        slide, Inches(0.18), top + Inches(0.05), Inches(12.8), Inches(0.62),
        text=title, font_size=22, bold=True, color=COLOR_DARK_BLUE,
    )


def _add_key_stat_callout(slide, stat_text: str):
    """슬라이드 우상단 대형 수치 callout (eval_section/case_study용)"""
    if not stat_text:
        return
    box = slide.shapes.add_shape(
        1,
        Inches(10.5), Inches(0.9), Inches(2.5), Inches(0.85),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_DARK_BLUE
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = stat_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = FONT_NAME
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_bullets_enhanced(slide, bullets: list, top=Inches(1.85), font_size=14,
                           highlight_phrase: str = ""):
    """향상된 bullet 목록 렌더링

    지원 형식:
    - str: 단순 bullet
    - dict {"text": "...", "sub": "..."}: 주요 메시지 + 들여쓰기 보조 설명

    highlight_phrase: 해당 문구를 포함한 run에 bold+COLOR_ACCENT 적용 (1개만)
    """
    if not bullets:
        return
    txb = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(5.2))
    tf = txb.text_frame
    tf.word_wrap = True

    highlighted = False  # 슬라이드 전체 1개만 강조

    for i, bullet in enumerate(bullets[:6]):
        if isinstance(bullet, dict):
            main_text = bullet.get("text", "")
            sub_text = bullet.get("sub", "")
        else:
            main_text = str(bullet)
            sub_text = ""

        # 주요 메시지 단락
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(2)
        p.line_spacing = 1.25

        # bullet 기호
        prefix_run = p.add_run()
        prefix_run.text = "• "
        prefix_run.font.name = FONT_NAME
        prefix_run.font.size = Pt(font_size)
        prefix_run.font.color.rgb = COLOR_ACCENT

        # 하이라이트 여부 판단
        should_highlight = (
            not highlighted
            and highlight_phrase
            and highlight_phrase in main_text
        )

        # 본문 run
        text_run = p.add_run()
        text_run.text = main_text
        text_run.font.name = FONT_NAME
        text_run.font.size = Pt(font_size)
        if should_highlight:
            text_run.font.bold = True
            text_run.font.color.rgb = COLOR_ACCENT
            highlighted = True
        else:
            text_run.font.color.rgb = COLOR_DARK_TEXT

        # sub-text 단락 (들여쓰기)
        if sub_text:
            sub_p = tf.add_paragraph()
            sub_p.space_before = Pt(1)
            sub_p.space_after = Pt(4)
            sub_p.line_spacing = 1.2
            sub_run = sub_p.add_run()
            sub_run.text = f"    └ {sub_text}"
            sub_run.font.name = FONT_NAME
            sub_run.font.size = Pt(11)
            sub_run.font.color.rgb = COLOR_SUB_TEXT


# ── 레이아웃별 렌더링 함수 ───────────────────────────────────────────────────

def _render_cover(slide, data: dict) -> None:
    """표지: 상단 다크블루 풀블리드 + 사업명 + Win Theme 부제목"""
    # 상단 1/3 다크블루 배경
    bg = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0), Inches(13.333), Inches(2.8),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DARK_BLUE
    bg.line.fill.background()

    # 사업명 (흰색 대형)
    _add_textbox(
        slide, Inches(1.2), Inches(0.7), Inches(11.0), Inches(1.2),
        text=data.get("title", ""),
        font_size=32, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
    )

    # 구분선
    line = slide.shapes.add_connector(
        1, Inches(1.5), Inches(3.1), Inches(11.8), Inches(3.1),
    )
    line.line.color.rgb = COLOR_ACCENT
    line.line.width = Pt(2)

    # Win Theme 부제목
    subtitle = data.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide, Inches(1.2), Inches(3.3), Inches(11.0), Inches(0.85),
            text=subtitle,
            font_size=18, bold=False, color=COLOR_ACCENT, align=PP_ALIGN.CENTER,
        )

    # highlight 문구 (있을 경우 대형 강조)
    highlight = data.get("highlight", {})
    if highlight and highlight.get("phrase"):
        _add_textbox(
            slide, Inches(1.2), Inches(4.3), Inches(11.0), Inches(0.6),
            text=highlight["phrase"],
            font_size=16, bold=True, color=COLOR_DARK_BLUE, align=PP_ALIGN.CENTER,
        )

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_key_message(slide, data: dict) -> None:
    """이기는 전략: Win Theme 헤드라인 + evidence bullets"""
    _add_slide_title(slide, data.get("title", "이기는 전략"))
    headline = data.get("headline", "")
    if headline:
        # 헤드라인 배경 박스
        hbg = slide.shapes.add_shape(
            1, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.72),
        )
        hbg.fill.solid()
        hbg.fill.fore_color.rgb = COLOR_DARK_BLUE
        hbg.line.fill.background()
        _add_textbox(
            slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.62),
            text=headline,
            font_size=16, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )

    highlight = data.get("highlight", {})
    _add_bullets_enhanced(
        slide, data.get("bullets", []),
        top=Inches(1.85),
        highlight_phrase=highlight.get("phrase", "") if highlight else "",
    )
    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_eval_section(slide, data: dict) -> None:
    """평가항목 섹션: eval_badge + 개선된 제목 + key_stat + 향상된 bullets"""
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", ""))
    _add_key_stat_callout(slide, data.get("key_stat", ""))

    highlight = data.get("highlight", {})
    _add_bullets_enhanced(
        slide, data.get("bullets", []),
        top=Inches(1.85),
        highlight_phrase=highlight.get("phrase", "") if highlight else "",
    )
    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_comparison(slide, data: dict) -> None:
    """차별화 비교: 경쟁사 vs 우리 — "우리" 열 강조"""
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", "차별화 포인트"))

    table_data = data.get("table", [])
    if not table_data:
        highlight = data.get("highlight", {})
        _add_bullets_enhanced(
            slide, data.get("bullets", []),
            highlight_phrase=highlight.get("phrase", "") if highlight else "",
        )
        _add_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    rows = len(table_data) + 1
    cols = 3
    tbl = slide.shapes.add_table(
        rows, cols, Inches(0.8), Inches(1.05), Inches(11.7), Inches(rows * 0.62),
    ).table

    competitor_label = data.get("competitor_label", "일반적 접근")
    for ci, header in enumerate(["구분", competitor_label, "우리 ✓"]):
        cell = tbl.cell(0, ci)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = FONT_NAME
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK_BLUE if ci != 2 else COLOR_ACCENT

    highlight_phrase = ""
    h = data.get("highlight", {})
    if h:
        highlight_phrase = h.get("phrase", "")

    for ri, row in enumerate(table_data[:6], start=1):
        for ci, (_, text) in enumerate([
            ("dimension", row.get("dimension", "")),
            ("competitor", row.get("competitor", "")),
            ("ours", row.get("ours", "")),
        ]):
            cell = tbl.cell(ri, ci)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = FONT_NAME
            run.font.size = Pt(13)
            if ci == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_OUR_BG
                run.font.bold = True
                run.font.color.rgb = COLOR_DARK_BLUE
                # 하이라이트 적용
                if highlight_phrase and highlight_phrase in text:
                    run.font.color.rgb = COLOR_ACCENT
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_timeline(slide, data: dict) -> None:
    """추진 계획 타임라인: 단계별 가로 배치 + 화살표 연결"""
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", "추진 계획"))

    phases = data.get("phases", []) or data.get("implementation_checklist", [])
    if not phases:
        _add_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    n = min(len(phases), 5)
    total_w = Inches(11.4)
    col_w = total_w / n
    block_w = col_w - Inches(0.2)
    start_left = Inches(0.9)
    top = Inches(1.05)
    height = Inches(5.6)

    phase_colors = [
        RGBColor(0x1F, 0x49, 0x7D),
        RGBColor(0x2E, 0x75, 0xB6),
        RGBColor(0x5B, 0x9B, 0xD5),
    ]

    for i, phase in enumerate(phases[:5]):
        left = start_left + col_w * i
        shape = slide.shapes.add_shape(1, left, top, block_w, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = phase_colors[i % len(phase_colors)]
        shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase.get("name", phase.get("phase", f"{i+1}단계"))
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = FONT_NAME
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        if phase.get("duration"):
            p2 = tf.add_paragraph()
            p2.text = phase["duration"]
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.runs[0] if p2.runs else p2.add_run()
            r2.font.name = FONT_NAME
            r2.font.size = Pt(11)
            r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        for d in phase.get("deliverables", [])[:3]:
            pd = tf.add_paragraph()
            pd.text = f"• {d}"
            pd.alignment = PP_ALIGN.LEFT
            rd = pd.runs[0] if pd.runs else pd.add_run()
            rd.font.name = FONT_NAME
            rd.font.size = Pt(10)
            rd.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        if i < n - 1:
            arrow_x = left + block_w + Inches(0.01)
            arrow_y = top + height / 2 - Inches(0.18)
            atxb = slide.shapes.add_textbox(arrow_x, arrow_y, Inches(0.18), Inches(0.36))
            ap = atxb.text_frame.paragraphs[0]
            ap.text = "▶"
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.runs[0] if ap.runs else ap.add_run()
            ar.font.size = Pt(11)
            ar.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_team(slide, data: dict) -> None:
    """투입 인력 구성 표"""
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", "투입 인력 & 역량"))

    team_rows = data.get("team_rows", [])
    if not team_rows:
        highlight = data.get("highlight", {})
        _add_bullets_enhanced(
            slide, data.get("bullets", []),
            highlight_phrase=highlight.get("phrase", "") if highlight else "",
        )
        _add_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    rows = len(team_rows) + 1
    tbl = slide.shapes.add_table(
        rows, 4, Inches(0.8), Inches(1.05), Inches(11.7), Inches(rows * 0.57),
    ).table

    for ci, header in enumerate(["역할", "등급", "투입 기간", "담당 업무"]):
        cell = tbl.cell(0, ci)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = FONT_NAME
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK_BLUE

    for ri, row in enumerate(team_rows, start=1):
        vals = [
            row.get("role", ""),
            row.get("grade", ""),
            str(row.get("person_months", "")) + "개월" if row.get("person_months") else row.get("duration", ""),
            row.get("task", ""),
        ]
        for ci, val in enumerate(vals):
            cell = tbl.cell(ri, ci)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 3 else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = FONT_NAME
            run.font.size = Pt(12)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_numbers_callout(slide, data: dict) -> None:
    """핵심 수치 강조: 대형 숫자 + 짧은 설명 카드 가로 배치"""
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", ""))

    numbers = data.get("numbers", [])
    if not numbers:
        highlight = data.get("highlight", {})
        _add_bullets_enhanced(
            slide, data.get("bullets", []),
            highlight_phrase=highlight.get("phrase", "") if highlight else "",
        )
        _add_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    n = min(len(numbers), 4)
    card_w = Inches(11.6) / n
    start_left = Inches(0.9)
    top = Inches(1.05)
    card_h = Inches(5.6)

    for i, item in enumerate(numbers[:4]):
        left = start_left + card_w * i
        card = slide.shapes.add_shape(1, left, top, card_w - Inches(0.15), card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xEA, 0xF1, 0xFB)
        card.line.color.rgb = COLOR_ACCENT
        card.line.width = Pt(1)

        value_top = top + Inches(0.6)
        _add_textbox(
            slide, left + Inches(0.1), value_top, card_w - Inches(0.2), Inches(1.3),
            text=str(item.get("value", "")),
            font_size=38, bold=True, color=COLOR_DARK_BLUE, align=PP_ALIGN.CENTER,
        )
        label_top = value_top + Inches(1.3)
        _add_textbox(
            slide, left + Inches(0.1), label_top, card_w - Inches(0.2), Inches(0.55),
            text=str(item.get("label", "")),
            font_size=13, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER,
        )
        desc_top = label_top + Inches(0.6)
        _add_textbox(
            slide, left + Inches(0.1), desc_top, card_w - Inches(0.2), Inches(1.8),
            text=str(item.get("description", "")),
            font_size=11, color=COLOR_DARK_TEXT, align=PP_ALIGN.CENTER,
        )

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_agenda(slide, data: dict) -> None:
    """목차/아젠다 슬라이드"""
    _add_slide_title(slide, data.get("title", "발표 순서"))

    items = data.get("items", [])
    if not items:
        _add_bullets_enhanced(slide, data.get("bullets", []))
        _add_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    top = Inches(1.05)
    item_h = Inches(0.68)

    for i, item in enumerate(items[:8]):
        item_top = top + item_h * i
        num_box = slide.shapes.add_shape(1, Inches(0.9), item_top, Inches(0.55), Inches(0.52))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_DARK_BLUE
        num_box.line.fill.background()
        ntf = num_box.text_frame
        np_ = ntf.paragraphs[0]
        np_.text = str(item.get("num", i + 1))
        np_.alignment = PP_ALIGN.CENTER
        nr = np_.runs[0] if np_.runs else np_.add_run()
        nr.font.name = FONT_NAME
        nr.font.size = Pt(13)
        nr.font.bold = True
        nr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        _add_textbox(
            slide, Inches(1.6), item_top + Inches(0.05), Inches(9.2), Inches(0.46),
            text=str(item.get("section", "")),
            font_size=14, bold=True, color=COLOR_DARK_TEXT,
        )
        score = item.get("score", "")
        if score:
            _add_textbox(
                slide, Inches(11.0), item_top + Inches(0.05), Inches(1.8), Inches(0.46),
                text=f"{score}점", font_size=13, bold=True,
                color=COLOR_ACCENT, align=PP_ALIGN.RIGHT,
            )
        if i < len(items) - 1:
            line_y = item_top + item_h - Inches(0.08)
            ln = slide.shapes.add_connector(1, Inches(0.9), line_y, Inches(12.4), line_y)
            ln.line.color.rgb = COLOR_LIGHT_GRAY
            ln.line.width = Pt(0.75)

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_process_flow(slide, data: dict) -> None:
    """프로세스 흐름도: 단계 박스 → 화살표"""
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", "추진 방법론"))

    steps = data.get("steps", [])
    if not steps:
        highlight = data.get("highlight", {})
        _add_bullets_enhanced(
            slide, data.get("bullets", []),
            highlight_phrase=highlight.get("phrase", "") if highlight else "",
        )
        _add_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    n = min(len(steps), 5)
    total_w = Inches(11.4)
    col_w = total_w / n
    box_w = col_w - Inches(0.3)
    start_left = Inches(0.9)
    top = Inches(1.05)
    box_h = Inches(5.6)

    step_colors = [
        RGBColor(0x1F, 0x49, 0x7D),
        RGBColor(0x2E, 0x75, 0xB6),
        RGBColor(0x5B, 0x9B, 0xD5),
    ]

    for i, step in enumerate(steps[:5]):
        left = start_left + col_w * i
        hdr = slide.shapes.add_shape(1, left, top, box_w, Inches(0.55))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = step_colors[i % len(step_colors)]
        hdr.line.fill.background()
        htf = hdr.text_frame
        hp = htf.paragraphs[0]
        hp.text = step.get("name", f"Step {i+1}")
        hp.alignment = PP_ALIGN.CENTER
        hr = hp.runs[0] if hp.runs else hp.add_run()
        hr.font.name = FONT_NAME
        hr.font.size = Pt(13)
        hr.font.bold = True
        hr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        body = slide.shapes.add_shape(1, left, top + Inches(0.55), box_w, box_h - Inches(0.55))
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFD)
        body.line.color.rgb = step_colors[i % len(step_colors)]
        body.line.width = Pt(1)

        desc = step.get("description", "")
        if desc:
            _add_textbox(
                slide, left + Inches(0.1), top + Inches(0.65), box_w - Inches(0.15), Inches(1.2),
                text=desc, font_size=12, color=COLOR_DARK_TEXT,
            )
        for j, out in enumerate(step.get("outputs", [])[:3]):
            out_top = top + Inches(1.95) + Inches(0.52) * j
            _add_textbox(
                slide, left + Inches(0.1), out_top, box_w - Inches(0.15), Inches(0.45),
                text=f"▸ {out}", font_size=11, color=COLOR_ACCENT,
            )

        if i < n - 1:
            ax = left + box_w + Inches(0.06)
            ay = top + box_h / 2 - Inches(0.18)
            atxb = slide.shapes.add_textbox(ax, ay, Inches(0.22), Inches(0.36))
            ap = atxb.text_frame.paragraphs[0]
            ap.text = "▶"
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.runs[0] if ap.runs else ap.add_run()
            ar.font.size = Pt(11)
            ar.font.color.rgb = COLOR_ACCENT

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_quote_highlight(slide, data: dict) -> None:
    """인용구 강조: 미션/비전/핵심 슬로건"""
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", ""))

    quote = data.get("quote", "")
    source = data.get("source", "")
    context = data.get("context", "")

    if not quote:
        highlight = data.get("highlight", {})
        _add_bullets_enhanced(
            slide, data.get("bullets", []),
            highlight_phrase=highlight.get("phrase", "") if highlight else "",
        )
        _add_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    bg = slide.shapes.add_shape(1, Inches(1.0), Inches(1.2), Inches(11.3), Inches(4.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xEA, 0xF1, 0xFB)
    bg.line.color.rgb = COLOR_ACCENT
    bg.line.width = Pt(1.5)

    _add_textbox(
        slide, Inches(1.2), Inches(1.15), Inches(1.0), Inches(1.0),
        text="\u201c",
        font_size=72, bold=True, color=COLOR_ACCENT,
    )
    _add_textbox(
        slide, Inches(2.0), Inches(1.6), Inches(9.5), Inches(3.0),
        text=quote,
        font_size=22, bold=True, color=COLOR_DARK_BLUE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, Inches(11.2), Inches(4.5), Inches(1.0), Inches(0.7),
        text="\u201d",
        font_size=72, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.RIGHT,
    )

    if source or context:
        combined = (
            f"— {source}" if source and not context
            else f"— {context}" if context and not source
            else f"— {source}  |  {context}"
        )
        _add_textbox(
            slide, Inches(1.5), Inches(5.9), Inches(10.3), Inches(0.5),
            text=combined,
            font_size=13, color=COLOR_ACCENT, align=PP_ALIGN.RIGHT,
        )

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_problem_sync(slide, data: dict) -> None:
    """문제 공감: 발주처 현황 통계 + 핵심 문제 + 당사 솔루션 연결"""
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", "발주처 현황과 핵심 과제"))

    problem = data.get("problem_statement", "")
    _cs = data.get("current_state", [])
    if isinstance(_cs, dict):
        stats_list = []
        metric = _cs.get("metric", "") or _cs.get("stat", "")
        label  = _cs.get("period", "") or _cs.get("label", "")
        if metric:
            stats_list.append({"stat": metric, "label": label})
        source = _cs.get("source", "")
        if source:
            stats_list.append({"stat": source, "label": ""})
        current_stats = stats_list
    else:
        current_stats = _cs if isinstance(_cs, list) else []
    pain_points = data.get("pain_points", [])

    left_bg = slide.shapes.add_shape(1, Inches(0.7), Inches(1.0), Inches(5.5), Inches(5.6))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = COLOR_DARK_BLUE
    left_bg.line.fill.background()

    _add_textbox(
        slide, Inches(0.9), Inches(1.1), Inches(5.1), Inches(0.5),
        text="현황 & 문제",
        font_size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
    )

    for i, stat_item in enumerate(current_stats[:3]):
        sy = Inches(1.75) + Inches(1.55) * i
        if isinstance(stat_item, dict):
            val = str(stat_item.get("stat", ""))
            lbl = str(stat_item.get("label", ""))
        else:
            val, lbl = str(stat_item), ""

        _add_textbox(
            slide, Inches(0.9), sy, Inches(5.0), Inches(0.85),
            text=val, font_size=30, bold=True,
            color=COLOR_YELLOW, align=PP_ALIGN.CENTER,
        )
        if lbl:
            _add_textbox(
                slide, Inches(0.9), sy + Inches(0.82), Inches(5.0), Inches(0.55),
                text=lbl, font_size=12,
                color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER,
            )

    right_top = Inches(1.0)
    _add_textbox(
        slide, Inches(6.6), right_top, Inches(6.1), Inches(0.5),
        text="핵심 과제",
        font_size=14, bold=True, color=COLOR_DARK_BLUE,
    )
    line = slide.shapes.add_connector(
        1, Inches(6.6), right_top + Inches(0.52), Inches(12.5), right_top + Inches(0.52)
    )
    line.line.color.rgb = COLOR_ACCENT
    line.line.width = Pt(1)

    for i, pt in enumerate(pain_points[:4]):
        pt_top = Inches(1.65) + Inches(1.05) * i
        num_box = slide.shapes.add_shape(1, Inches(6.6), pt_top, Inches(0.42), Inches(0.4))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_ACCENT
        num_box.line.fill.background()
        ntf = num_box.text_frame
        np_ = ntf.paragraphs[0]
        np_.text = str(i + 1)
        np_.alignment = PP_ALIGN.CENTER
        nr = np_.runs[0] if np_.runs else np_.add_run()
        nr.font.name = FONT_NAME
        nr.font.size = Pt(11)
        nr.font.bold = True
        nr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        _add_textbox(
            slide, Inches(7.15), pt_top, Inches(5.4), Inches(0.6),
            text=str(pt), font_size=13, color=COLOR_DARK_TEXT,
        )

    if problem:
        sol_bg = slide.shapes.add_shape(1, Inches(6.6), Inches(6.05), Inches(6.1), Inches(0.6))
        sol_bg.fill.solid()
        sol_bg.fill.fore_color.rgb = RGBColor(0xD6, 0xE4, 0xF0)
        sol_bg.line.color.rgb = COLOR_ACCENT
        sol_bg.line.width = Pt(1)
        _add_textbox(
            slide, Inches(6.7), Inches(6.1), Inches(5.9), Inches(0.5),
            text=f"▶ {problem}",
            font_size=12, bold=True, color=COLOR_DARK_BLUE,
        )

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_closing(slide, data: dict) -> None:
    """마무리: Win Theme 재강조 + 차별화 포인트"""
    _add_slide_title(slide, data.get("title", "왜 우리인가"))
    headline = data.get("headline", "")
    if headline:
        hbg = slide.shapes.add_shape(
            1, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.72),
        )
        hbg.fill.solid()
        hbg.fill.fore_color.rgb = COLOR_DARK_BLUE
        hbg.line.fill.background()
        _add_textbox(
            slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.62),
            text=headline,
            font_size=16, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )

    highlight = data.get("highlight", {})
    _add_bullets_enhanced(
        slide, data.get("bullets", []),
        top=Inches(1.85), font_size=16,
        highlight_phrase=highlight.get("phrase", "") if highlight else "",
    )
    _add_speaker_notes(slide, data.get("speaker_notes", ""))


# ── 신규 레이아웃 ─────────────────────────────────────────────────────────────

def _render_split_panel(slide, data: dict) -> None:
    """F-Pattern 분할 패널: 좌측 차트(55%) + 우측 번호형 전략 포인트(45%)

    chart_data 구조: {chart_type: "bar|line", categories: [...], series: [{name, values}]}
    f_pattern_anchor: 차트 하단 대형 강조 수치 (F-Pattern 첫 시선)
    points: [{headline, sub_text, emphasis_value}]
    """
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", ""))

    # ── 좌측 패널 (차트) ──────────────────────────────────────────────────
    chart_data_raw = data.get("chart_data", {})
    chart_title = data.get("chart_title", "")
    left_area_left = Inches(0.5)
    left_area_top = Inches(1.05)
    left_area_w = Inches(7.0)
    left_area_h = Inches(5.6)

    # 차트 제목 (결론형 문장)
    if chart_title:
        _add_textbox(
            slide, left_area_left, left_area_top, left_area_w, Inches(0.55),
            text=chart_title,
            font_size=12, bold=True, color=COLOR_DARK_BLUE,
        )

    # 차트 렌더링
    if chart_data_raw and chart_data_raw.get("categories"):
        try:
            cd = ChartData()
            cd.categories = chart_data_raw.get("categories", [])
            for series in chart_data_raw.get("series", []):
                cd.add_series(
                    series.get("name", ""),
                    tuple(float(v) for v in series.get("values", [])),
                )
            chart_type_str = chart_data_raw.get("chart_type", "bar").lower()
            xl_type = (
                XL_CHART_TYPE.LINE
                if chart_type_str == "line"
                else XL_CHART_TYPE.COLUMN_CLUSTERED
            )
            chart_top = left_area_top + (Inches(0.6) if chart_title else Inches(0))
            chart_h = left_area_h - (Inches(0.6) if chart_title else Inches(0)) - Inches(0.9)
            chart_shape = slide.shapes.add_chart(
                xl_type,
                left_area_left, chart_top, left_area_w, chart_h,
                cd,
            )
            chart = chart_shape.chart
            # 차트 스타일 정리 (불필요한 요소 제거)
            chart.has_legend = len(chart_data_raw.get("series", [])) > 1
            if chart.has_title:
                chart.chart_title.has_text_frame = False
            # 첫 번째 시리즈 색상 강조
            if chart.series:
                for s_idx, series_obj in enumerate(chart.series):
                    try:
                        pt = series_obj.points[len(series_obj.points) - 1]
                        pt.format.fill.solid()
                        pt.format.fill.fore_color.rgb = COLOR_DARK_BLUE if s_idx == 0 else COLOR_ACCENT
                    except Exception:
                        pass
        except Exception as e:
            logger.warning(f"차트 렌더링 실패 (split_panel), 박스로 대체: {e}")
            chart_box = slide.shapes.add_shape(
                1, left_area_left, left_area_top + Inches(0.6),
                left_area_w, left_area_h - Inches(0.9),
            )
            chart_box.fill.solid()
            chart_box.fill.fore_color.rgb = COLOR_TITLE_BG
            chart_box.line.color.rgb = COLOR_ACCENT
            chart_box.line.width = Pt(1)
            _add_textbox(
                slide,
                left_area_left + Inches(0.2),
                left_area_top + left_area_h / 2 - Inches(0.3),
                left_area_w - Inches(0.4), Inches(0.6),
                text="[차트 데이터 없음]",
                font_size=12, color=COLOR_SUB_TEXT, align=PP_ALIGN.CENTER,
            )

    # F-Pattern anchor: 차트 하단 대형 강조 수치
    anchor = data.get("f_pattern_anchor", "")
    if anchor:
        anchor_top = left_area_top + left_area_h - Inches(0.75)
        anchor_bg = slide.shapes.add_shape(
            1, left_area_left, anchor_top, left_area_w, Inches(0.65),
        )
        anchor_bg.fill.solid()
        anchor_bg.fill.fore_color.rgb = COLOR_DARK_BLUE
        anchor_bg.line.fill.background()
        _add_textbox(
            slide, left_area_left + Inches(0.1), anchor_top + Inches(0.05),
            left_area_w - Inches(0.2), Inches(0.55),
            text=anchor,
            font_size=26, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )

    # ── 우측 패널 (번호형 포인트) ─────────────────────────────────────────
    right_left = Inches(7.8)
    right_top = Inches(1.05)
    right_w = Inches(5.0)
    right_h = Inches(5.6)

    points = data.get("points", [])
    if not points:
        # fallback: bullets 사용
        highlight = data.get("highlight", {})
        _add_bullets_enhanced(
            slide, data.get("bullets", []),
            top=right_top,
            highlight_phrase=highlight.get("phrase", "") if highlight else "",
        )
    else:
        n_pts = min(len(points), 4)
        pt_h = right_h / n_pts

        highlight_phrase = ""
        h = data.get("highlight", {})
        if h:
            highlight_phrase = h.get("phrase", "")
        highlighted = False

        for i, pt in enumerate(points[:4]):
            pt_top = right_top + pt_h * i

            # 번호 원형
            num_circle = slide.shapes.add_shape(
                1, right_left, pt_top + Inches(0.1), Inches(0.45), Inches(0.42),
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = COLOR_DARK_BLUE
            num_circle.line.fill.background()
            ntf = num_circle.text_frame
            np_ = ntf.paragraphs[0]
            np_.text = str(i + 1)
            np_.alignment = PP_ALIGN.CENTER
            nr = np_.runs[0] if np_.runs else np_.add_run()
            nr.font.name = FONT_NAME
            nr.font.size = Pt(12)
            nr.font.bold = True
            nr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # 헤드라인 텍스트박스
            headline_text = pt.get("headline", "")
            should_hl = (
                not highlighted
                and highlight_phrase
                and highlight_phrase in headline_text
            )
            h_txb = slide.shapes.add_textbox(
                right_left + Inches(0.55), pt_top + Inches(0.05),
                right_w - Inches(0.65), Inches(0.48),
            )
            htf = h_txb.text_frame
            hp = htf.paragraphs[0]
            hr = hp.add_run()
            hr.text = headline_text
            hr.font.name = FONT_NAME
            hr.font.size = Pt(13)
            hr.font.bold = True
            hr.font.color.rgb = COLOR_ACCENT if should_hl else COLOR_DARK_BLUE
            if should_hl:
                highlighted = True

            # sub_text
            sub = pt.get("sub_text", "")
            if sub:
                _add_textbox(
                    slide,
                    right_left + Inches(0.55), pt_top + Inches(0.55),
                    right_w - Inches(0.65), Inches(0.5),
                    text=sub, font_size=11, color=COLOR_SUB_TEXT,
                )

            # emphasis_value (우측 강조 수치)
            emph = pt.get("emphasis_value", "")
            if emph:
                _add_textbox(
                    slide,
                    right_left + Inches(0.55), pt_top + Inches(0.95),
                    right_w - Inches(0.65), Inches(0.4),
                    text=emph, font_size=12, bold=True, color=COLOR_ACCENT,
                )

            # 구분선 (마지막 제외)
            if i < n_pts - 1:
                line_y = pt_top + pt_h - Inches(0.08)
                ln = slide.shapes.add_connector(
                    1, right_left, line_y, right_left + right_w - Inches(0.1), line_y,
                )
                ln.line.color.rgb = COLOR_LIGHT_GRAY
                ln.line.width = Pt(0.75)

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_numbered_strategy(slide, data: dict) -> None:
    """번호형 전략 포인트: 아이콘 원형 + bold 헤드라인 + sub-description + 강조 수치

    points: [{num, headline, sub_text, emphasis_value}]
    """
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", ""))

    points = data.get("points", [])
    if not points:
        highlight = data.get("highlight", {})
        _add_bullets_enhanced(
            slide, data.get("bullets", []),
            highlight_phrase=highlight.get("phrase", "") if highlight else "",
        )
        _add_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    n = min(len(points), 4)
    start_top = Inches(1.1)
    total_h = Inches(5.5)
    pt_h = total_h / n

    highlight_phrase = ""
    h = data.get("highlight", {})
    if h:
        highlight_phrase = h.get("phrase", "")
    highlighted = False

    for i, pt in enumerate(points[:4]):
        pt_top = start_top + pt_h * i

        # 번호 원형 (좌측)
        num_circle = slide.shapes.add_shape(
            1, Inches(0.5), pt_top + Inches(0.08), Inches(0.52), Inches(0.5),
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLOR_DARK_BLUE
        num_circle.line.fill.background()
        ntf = num_circle.text_frame
        np_ = ntf.paragraphs[0]
        np_.text = str(pt.get("num", i + 1))
        np_.alignment = PP_ALIGN.CENTER
        nr = np_.runs[0] if np_.runs else np_.add_run()
        nr.font.name = FONT_NAME
        nr.font.size = Pt(14)
        nr.font.bold = True
        nr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 헤드라인
        headline_text = pt.get("headline", "")
        should_hl = (
            not highlighted
            and highlight_phrase
            and highlight_phrase in headline_text
        )
        h_txb = slide.shapes.add_textbox(
            Inches(1.2), pt_top, Inches(9.5), Inches(0.52),
        )
        htf = h_txb.text_frame
        hp = htf.paragraphs[0]
        hr = hp.add_run()
        hr.text = headline_text
        hr.font.name = FONT_NAME
        hr.font.size = Pt(15)
        hr.font.bold = True
        hr.font.color.rgb = COLOR_ACCENT if should_hl else COLOR_DARK_BLUE
        if should_hl:
            highlighted = True

        # emphasis_value (우측 강조)
        emph = pt.get("emphasis_value", "")
        if emph:
            _add_textbox(
                slide, Inches(10.8), pt_top, Inches(2.2), Inches(0.52),
                text=emph, font_size=15, bold=True, color=COLOR_ACCENT,
                align=PP_ALIGN.RIGHT,
            )

        # sub_text
        sub = pt.get("sub_text", "")
        if sub:
            _add_textbox(
                slide, Inches(1.2), pt_top + Inches(0.52), Inches(11.5), Inches(0.52),
                text=f"└ {sub}", font_size=12, color=COLOR_SUB_TEXT,
            )

        # 구분선 (마지막 제외)
        if i < n - 1:
            line_y = pt_top + pt_h - Inches(0.06)
            ln = slide.shapes.add_connector(
                1, Inches(0.5), line_y, Inches(12.5), line_y,
            )
            ln.line.color.rgb = COLOR_LIGHT_GRAY
            ln.line.width = Pt(0.75)

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


def _render_case_study(slide, data: dict) -> None:
    """유사 사례/레퍼런스 증거 슬라이드

    cases: [{client, project, challenge, solution, result, relevance}]
    2~3개 사례를 카드 형태로 가로 배치, key_stat 우상단 대형 callout
    """
    _add_eval_badge(slide, data.get("eval_badge", ""))
    _add_slide_title(slide, data.get("title", "수행 사례"))
    _add_key_stat_callout(slide, data.get("key_stat", ""))

    cases = data.get("cases", [])
    if not cases:
        highlight = data.get("highlight", {})
        _add_bullets_enhanced(
            slide, data.get("bullets", []),
            highlight_phrase=highlight.get("phrase", "") if highlight else "",
        )
        _add_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    n = min(len(cases), 3)
    total_w = Inches(12.3)
    card_w = total_w / n
    start_left = Inches(0.5)
    top = Inches(1.1)
    card_h = Inches(5.55)

    highlight_phrase = ""
    h = data.get("highlight", {})
    if h:
        highlight_phrase = h.get("phrase", "")

    for i, case in enumerate(cases[:3]):
        left = start_left + card_w * i
        cw = card_w - Inches(0.15)

        # 카드 배경
        card_bg = slide.shapes.add_shape(1, left, top, cw, card_h)
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLOR_CASE_BG
        card_bg.line.color.rgb = COLOR_ACCENT
        card_bg.line.width = Pt(1)

        # 헤더 (기관명)
        hdr = slide.shapes.add_shape(1, left, top, cw, Inches(0.48))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = COLOR_DARK_BLUE
        hdr.line.fill.background()
        htf = hdr.text_frame
        hp = htf.paragraphs[0]
        hp.text = case.get("client", f"사례 {i+1}")
        hp.alignment = PP_ALIGN.CENTER
        hr = hp.runs[0] if hp.runs else hp.add_run()
        hr.font.name = FONT_NAME
        hr.font.size = Pt(13)
        hr.font.bold = True
        hr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 사업명
        _add_textbox(
            slide, left + Inches(0.1), top + Inches(0.55), cw - Inches(0.2), Inches(0.5),
            text=case.get("project", ""),
            font_size=12, bold=True, color=COLOR_DARK_BLUE,
        )

        # 항목 렌더링 함수
        def _case_row(label: str, value: str, y_offset: float, is_result: bool = False):
            row_top = top + Inches(y_offset)
            label_txb = slide.shapes.add_textbox(
                left + Inches(0.1), row_top, Inches(0.9), Inches(0.42),
            )
            ltf = label_txb.text_frame
            lp = ltf.paragraphs[0]
            lr = lp.add_run()
            lr.text = label
            lr.font.name = FONT_NAME
            lr.font.size = Pt(10)
            lr.font.bold = True
            lr.font.color.rgb = COLOR_ACCENT

            val_txb = slide.shapes.add_textbox(
                left + Inches(1.05), row_top, cw - Inches(1.2), Inches(0.55),
            )
            vtf = val_txb.text_frame
            vtf.word_wrap = True
            vp = vtf.paragraphs[0]
            vr = vp.add_run()
            vr.text = value
            vr.font.name = FONT_NAME
            vr.font.size = Pt(11)
            should_bold = is_result and highlight_phrase and highlight_phrase in value
            vr.font.bold = should_bold or is_result
            vr.font.color.rgb = COLOR_ACCENT if should_bold else (
                COLOR_DARK_BLUE if is_result else COLOR_DARK_TEXT
            )

        _case_row("과제", case.get("challenge", ""), 1.12)
        _case_row("해결", case.get("solution", ""), 1.72)
        _case_row("성과", case.get("result", ""), 2.32, is_result=True)
        _case_row("연관성", case.get("relevance", ""), 3.1)

    _add_speaker_notes(slide, data.get("speaker_notes", ""))


# ── 슬라이드 디스패처 ────────────────────────────────────────────────────────

_LAYOUT_MAP = {
    "cover":               _render_cover,
    "key_message":         _render_key_message,
    "eval_section":        _render_eval_section,
    "comparison":          _render_comparison,
    "timeline":            _render_timeline,
    "team":                _render_team,
    "closing":             _render_closing,
    "numbers_callout":     _render_numbers_callout,
    "agenda":              _render_agenda,
    "process_flow":        _render_process_flow,
    "quote_highlight":     _render_quote_highlight,
    "problem_sync":        _render_problem_sync,
    # 신규
    "split_panel":         _render_split_panel,
    "numbered_strategy":   _render_numbered_strategy,
    "case_study":          _render_case_study,
}


def _render_slide(prs: Presentation, data: dict) -> None:
    """레이아웃 타입별 슬라이드 렌더링 (실패 시 텍스트 fallback)"""
    layout = data.get("layout", "eval_section")
    renderer = _LAYOUT_MAP.get(layout, _render_eval_section)
    try:
        layout_idx = min(6, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        renderer(slide, data)
        if layout != "cover":
            _add_slide_number(slide, data.get("slide_num", 0))
    except Exception as e:
        logger.warning(f"슬라이드 렌더링 실패 ({layout}), fallback 적용: {e}")
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = data.get("title", "")
            bullets = data.get("bullets", [])
            if bullets and len(slide.placeholders) > 1:
                slide.placeholders[1].text = "\n".join(
                    b if isinstance(b, str) else b.get("text", "") for b in bullets
                )
        except Exception as e2:
            logger.error(f"fallback 슬라이드도 실패: {e2}")


# ── 템플릿 초기화 ─────────────────────────────────────────────────────────────

def _init_presentation(template_path: Optional[Path]) -> Presentation:
    """템플릿 파일 로드 또는 빈 프레젠테이션 생성"""
    if template_path and template_path.exists():
        try:
            prs = Presentation(str(template_path))
            while len(prs.slides._sldIdLst) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
            logger.info(f"템플릿 로드 완료: {template_path}")
            return prs
        except Exception as e:
            logger.warning(f"템플릿 로드 실패, scratch로 fallback: {e}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


# ── 공개 인터페이스 ───────────────────────────────────────────────────────────

def build_presentation_pptx(
    slides_json: dict,
    output_path: Path,
    project_name: str = "",
    template_path: Optional[Path] = None,
) -> Path:
    """
    슬라이드 JSON → 발표용 PPTX 파일 생성

    Args:
        slides_json: presentation_generator.py가 반환한 슬라이드 JSON
        output_path: 저장 경로
        project_name: 표지 fallback 제목
        template_path: PPTX 템플릿 파일 경로 (None이면 scratch)

    Returns:
        저장된 PPTX 파일 경로
    """
    prs = _init_presentation(template_path)

    slides = slides_json.get("slides", [])
    if not slides:
        logger.warning("슬라이드 데이터 없음 — 빈 표지만 생성")
        slides = [{"slide_num": 1, "layout": "cover", "title": project_name}]

    for slide_data in sorted(slides, key=lambda s: s.get("slide_num", 99)):
        _render_slide(prs, slide_data)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info(f"PPTX 저장 완료: {output_path} ({len(slides)}장)")
    return output_path
