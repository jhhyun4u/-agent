from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def build_pptx(sections: dict, output_path: Path, project_name: str = "용역 제안서") -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 표지
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{project_name}"
    slide.placeholders[1].text = "용역 제안서"

    for heading, body in sections.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = heading
        tf = slide.placeholders[1].text_frame
        tf.clear()
        if isinstance(body, str):
            lines = [line.strip() for line in body.split("\n") if line.strip()]
            for i, line in enumerate(lines[:10]):  # 슬라이드당 최대 10줄
                if i == 0:
                    tf.paragraphs[0].text = line
                    tf.paragraphs[0].font.size = Pt(16)
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(16)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
