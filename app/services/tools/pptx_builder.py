"""
PPTX 빌더 (§12-4)

그래프 상태의 ppt_slides → PowerPoint 바이트 변환.
- 표지 + 목차 + 본문 슬라이드 + 마무리
- presentation_strategy 반영 (키 메시지, 발표 시간 배분)
"""

import asyncio
import io
import logging
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# ── 레이아웃 상수 ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TITLE_FONT_SIZE = Pt(28)
SUBTITLE_FONT_SIZE = Pt(18)
BODY_FONT_SIZE = Pt(16)
SMALL_FONT_SIZE = Pt(14)
MAX_LINES_PER_SLIDE = 12


async def build_pptx(
    slides: list[dict[str, Any]],
    proposal_name: str = "용역 제안서",
    presentation_strategy: dict[str, Any] | None = None,
) -> bytes:
    """ppt_slides 리스트 → PPTX 바이트."""
    return await asyncio.to_thread(
        _build_pptx_sync, slides, proposal_name, presentation_strategy,
    )


def _build_pptx_sync(
    slides: list[dict[str, Any]],
    proposal_name: str,
    presentation_strategy: dict[str, Any] | None,
) -> bytes:
    """동기 PPTX 빌드 (to_thread용)."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _add_title_slide(prs, proposal_name, presentation_strategy)
    _add_toc_slide(prs, slides)

    for slide_data in slides:
        _add_content_slide(prs, slide_data)

    _add_closing_slide(prs, proposal_name)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _add_title_slide(
    prs: Presentation,
    proposal_name: str,
    strategy: dict[str, Any] | None,
):
    """표지 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = proposal_name
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = TITLE_FONT_SIZE
            p.alignment = PP_ALIGN.CENTER

    # 부제: 키 메시지 또는 기본 텍스트
    subtitle = "용역 제안서"
    if strategy and strategy.get("key_message"):
        subtitle = strategy["key_message"]

    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
        for p in slide.placeholders[1].text_frame.paragraphs:
            p.font.size = SUBTITLE_FONT_SIZE
            p.alignment = PP_ALIGN.CENTER


def _add_toc_slide(prs: Presentation, slides: list[dict[str, Any]]):
    """목차 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    if slide.shapes.title:
        slide.shapes.title.text = "목 차"
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE

    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, s in enumerate(slides):
            title = s.get("title", s.get("section_title", f"슬라이드 {i + 1}"))
            if i == 0:
                tf.paragraphs[0].text = f"{i + 1}. {title}"
                tf.paragraphs[0].font.size = BODY_FONT_SIZE
            else:
                p = tf.add_paragraph()
                p.text = f"{i + 1}. {title}"
                p.font.size = BODY_FONT_SIZE


def _add_content_slide(prs: Presentation, slide_data: dict[str, Any]):
    """본문 슬라이드 1장."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_data.get("title", slide_data.get("section_title", ""))

    if slide.shapes.title:
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE

    # 본문 내용 추출
    body = slide_data.get("content", slide_data.get("body", ""))
    bullets = slide_data.get("bullets", [])
    notes = slide_data.get("speaker_notes", "")

    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()

        if bullets:
            _render_bullets(tf, bullets)
        elif body:
            _render_body(tf, body)

    # 발표자 노트
    if notes and slide.has_notes_slide:
        slide.notes_slide.notes_text_frame.text = notes


def _render_bullets(tf, bullets: list[str]):
    """불릿 리스트 렌더링."""
    for i, bullet in enumerate(bullets[:MAX_LINES_PER_SLIDE]):
        if i == 0:
            tf.paragraphs[0].text = bullet
            tf.paragraphs[0].font.size = BODY_FONT_SIZE
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = BODY_FONT_SIZE


def _render_body(tf, body: str):
    """텍스트 본문 렌더링 (줄바꿈 기준 분할)."""
    lines = [line.strip() for line in body.split("\n") if line.strip()]
    for i, line in enumerate(lines[:MAX_LINES_PER_SLIDE]):
        if i == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = BODY_FONT_SIZE
        else:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = BODY_FONT_SIZE


def _add_closing_slide(prs: Presentation, proposal_name: str):
    """마무리 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "감사합니다"
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = TITLE_FONT_SIZE
            p.alignment = PP_ALIGN.CENTER

    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = proposal_name
        for p in slide.placeholders[1].text_frame.paragraphs:
            p.font.size = SUBTITLE_FONT_SIZE
            p.alignment = PP_ALIGN.CENTER


# ── 레거시 호환 ──

def build_pptx_legacy(sections: dict, output_path, project_name: str = "용역 제안서"):
    """기존 동기 인터페이스 (하위 호환)."""
    from pathlib import Path
    output_path = Path(output_path)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = project_name
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "용역 제안서"

    for heading, body in sections.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = heading
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            if isinstance(body, str):
                lines = [line.strip() for line in body.split("\n") if line.strip()]
                for i, line in enumerate(lines[:10]):
                    if i == 0:
                        tf.paragraphs[0].text = line
                        tf.paragraphs[0].font.size = BODY_FONT_SIZE
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = BODY_FONT_SIZE

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
